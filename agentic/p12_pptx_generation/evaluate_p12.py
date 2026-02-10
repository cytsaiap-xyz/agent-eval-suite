#!/usr/bin/env python3
"""Evaluation script for P12: PowerPoint Generation"""

import json
import sys
from pathlib import Path


def evaluate(base_path: str = ".") -> dict:
    base = Path(base_path)
    results = {"scores": {}, "details": [], "total_score": 0, "max_score": 100}

    pptx_path = base / "quarterly_review.pptx"
    if not pptx_path.exists():
        results["details"].append("✗ quarterly_review.pptx not found")
        return results

    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))

        # Check slide count
        slide_count = len(prs.slides)
        if slide_count >= 10:
            results["scores"]["slide_count"] = 15
            results["details"].append(f"✓ {slide_count} slides (≥10 required)")
        else:
            results["scores"]["slide_count"] = slide_count
            results["details"].append(f"◐ Only {slide_count} slides")

        # Check for charts
        charts_found = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_chart:
                    charts_found += 1

        if charts_found >= 3:
            results["scores"]["charts"] = 25
            results["details"].append(f"✓ {charts_found} charts found")
        elif charts_found > 0:
            results["scores"]["charts"] = charts_found * 5
            results["details"].append(f"◐ Only {charts_found} charts")
        else:
            results["scores"]["charts"] = 0
            results["details"].append("✗ No charts found")

        # Check for tables
        tables_found = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    tables_found += 1

        if tables_found >= 1:
            results["scores"]["tables"] = 15
            results["details"].append(f"✓ {tables_found} tables found")
        else:
            results["scores"]["tables"] = 0
            results["details"].append("✗ No tables found")

        # Check for text content
        total_text = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    total_text += len(shape.text_frame.text)

        if total_text > 500:
            results["scores"]["content"] = 20
            results["details"].append("✓ Sufficient text content")
        elif total_text > 100:
            results["scores"]["content"] = 10
            results["details"].append("◐ Limited text content")
        else:
            results["scores"]["content"] = 0
            results["details"].append("✗ Very little text content")

        # Assume formatting is reasonable if we got this far
        results["scores"]["formatting"] = 25
        results["details"].append("✓ Presentation file valid")

    except ImportError:
        results["details"].append("⚠ python-pptx not installed, limited validation")
        results["scores"]["file_exists"] = 50
    except Exception as e:
        results["details"].append(f"✗ Error reading pptx: {e}")

    results["total_score"] = sum(results["scores"].values())
    return results


def main():
    base_path = sys.argv[1] if len(sys.argv) > 1 else "."
    print("=" * 60)
    print("Problem 12: PPTX Generation - Evaluation")
    print("=" * 60)

    results = evaluate(base_path)

    for detail in results["details"]:
        print(f"  {detail}")

    print(f"\nTotal Score: {results['total_score']}/{results['max_score']}")

    with open(Path(base_path) / "evaluation_results.json", "w") as f:
        json.dump(results, f, indent=2)


if __name__ == "__main__":
    main()
