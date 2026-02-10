#!/usr/bin/env python3
"""
Generate all input files for P15 Multi-Format Integration.
Requires: pip install openpyxl python-docx python-pptx reportlab
"""

from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Table, TableStyle, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from datetime import datetime
import random


def create_sales_excel():
    """Create sales_report.xlsx"""
    wb = Workbook()

    # Summary sheet
    ws_summary = wb.active
    ws_summary.title = "Summary"
    ws_summary.append(["Q4 2024 Sales Summary"])
    ws_summary.append([])
    ws_summary.append(["Metric", "Value"])
    ws_summary.append(["Total Revenue", 14400000])
    ws_summary.append(["Total Orders", 2680])
    ws_summary.append(["Average Order Value", 5373])
    ws_summary.append(["New Customers", 450])
    ws_summary.append(["Repeat Rate", "68%"])

    # Monthly data
    ws_monthly = wb.create_sheet("Monthly")
    ws_monthly.append(["Month", "Revenue", "Orders", "Customers"])
    monthly_data = [
        ["October", 4200000, 820, 650],
        ["November", 4800000, 900, 720],
        ["December", 5400000, 960, 810],
    ]
    for row in monthly_data:
        ws_monthly.append(row)

    # Product data
    ws_products = wb.create_sheet("Products")
    ws_products.append(["Product", "Revenue", "Units", "Margin"])
    product_data = [
        ["Product X", 4200000, 28000, "45%"],
        ["Product Y", 5800000, 42000, "38%"],
        ["Product Z", 3100000, 18000, "52%"],
        ["Services", 1300000, None, "65%"],
    ]
    for row in product_data:
        ws_products.append(row)

    # Regional data
    ws_regions = wb.create_sheet("Regions")
    ws_regions.append(["Region", "Revenue", "Growth"])
    regional_data = [
        ["North America", 7200000, "12%"],
        ["Europe", 4300000, "15%"],
        ["APAC", 2900000, "45%"],
    ]
    for row in regional_data:
        ws_regions.append(row)

    wb.save("sales_report.xlsx")
    print("Created sales_report.xlsx")


def create_marketing_pptx():
    """Create marketing_deck.pptx"""
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Q4 2024 Marketing Review"

    # Campaign performance slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
    txBox.text_frame.text = "Campaign Performance"

    # Add table
    rows, cols = 5, 4
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(2)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    headers = ["Campaign", "Spend", "Leads", "CPL"]
    data = [
        ["Digital Ads", "$250,000", "5,200", "$48"],
        ["Content Marketing", "$80,000", "2,100", "$38"],
        ["Events", "$120,000", "1,800", "$67"],
        ["Email", "$30,000", "1,500", "$20"],
    ]

    for i, header in enumerate(headers):
        table.cell(0, i).text = header

    for row_idx, row_data in enumerate(data, 1):
        for col_idx, cell_data in enumerate(row_data):
            table.cell(row_idx, col_idx).text = cell_data

    # Market share slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
    txBox.text_frame.text = "Market Share Analysis"

    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
    tf = txBox2.text_frame
    tf.text = "Market Share: 18% (+2% YoY)\n"
    tf.text += "Brand Awareness: 45% (+8% YoY)\n"
    tf.text += "Customer Satisfaction: 4.2/5.0\n"
    tf.text += "Social Media Followers: 125,000 (+35%)"

    prs.save("marketing_deck.pptx")
    print("Created marketing_deck.pptx")


def create_finance_pdf():
    """Create finance_summary.pdf"""
    doc = SimpleDocTemplate("finance_summary.pdf", pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    story.append(Paragraph("Q4 2024 Financial Summary", styles['Heading1']))
    story.append(Spacer(1, 20))

    # Income statement
    story.append(Paragraph("Income Statement", styles['Heading2']))
    income_data = [
        ['Item', 'Q4 2024', 'Q4 2023', 'Change'],
        ['Revenue', '$14.4M', '$12.5M', '+15.2%'],
        ['Cost of Goods Sold', '$8.1M', '$7.2M', '+12.5%'],
        ['Gross Profit', '$6.3M', '$5.3M', '+18.9%'],
        ['Operating Expenses', '$3.5M', '$3.2M', '+9.4%'],
        ['Operating Income', '$2.8M', '$2.1M', '+33.3%'],
        ['Net Income', '$2.0M', '$1.5M', '+33.3%'],
    ]

    t = Table(income_data, colWidths=[150, 80, 80, 80])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1E3A5F')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (1, 0), (-1, -1), 'RIGHT'),
        ('GRID', (0, 0), (-1, -1), 1, colors.grey),
    ]))
    story.append(t)
    story.append(Spacer(1, 20))

    # Key ratios
    story.append(Paragraph("Key Financial Ratios", styles['Heading2']))
    story.append(Paragraph("""
    Gross Margin: 43.8% (vs 42.4% prior year)
    Operating Margin: 19.4% (vs 16.8% prior year)
    Net Margin: 13.9% (vs 12.0% prior year)
    Return on Assets: 15.4%
    Return on Equity: 24.1%
    Current Ratio: 2.8x
    Debt to Equity: 0.24x
    """, styles['Normal']))

    doc.build(story)
    print("Created finance_summary.pdf")


def create_operations_docx():
    """Create operations_update.docx"""
    doc = Document()

    doc.add_heading("Q4 2024 Operations Update", 0)

    doc.add_heading("Executive Summary", level=1)
    doc.add_paragraph("""
    Operations delivered strong performance in Q4 2024, achieving all key targets
    while supporting the company's growth initiatives. Manufacturing efficiency
    improved by 8%, and we successfully launched three new product lines.
    """)

    doc.add_heading("Key Performance Indicators", level=1)

    table = doc.add_table(rows=6, cols=4)
    table.style = 'Table Grid'

    headers = ["KPI", "Target", "Actual", "Status"]
    header_cells = table.rows[0].cells
    for i, header in enumerate(headers):
        header_cells[i].text = header

    kpi_data = [
        ["Production Volume", "150,000", "158,000", "Exceeded"],
        ["Defect Rate", "<1.0%", "0.5%", "Exceeded"],
        ["On-Time Delivery", ">95%", "98%", "Exceeded"],
        ["Inventory Turns", "8x", "8.2x", "Met"],
        ["Safety Incidents", "0", "0", "Met"],
    ]

    for row_idx, row_data in enumerate(kpi_data, 1):
        row_cells = table.rows[row_idx].cells
        for col_idx, cell_data in enumerate(row_data):
            row_cells[col_idx].text = cell_data

    doc.add_heading("Initiative Updates", level=1)
    doc.add_paragraph("Automation Project: Phase 2 complete, 15% efficiency gain", style='List Bullet')
    doc.add_paragraph("New Warehouse: Construction on schedule, Q2 2025 opening", style='List Bullet')
    doc.add_paragraph("Supplier Diversification: Added 3 new qualified suppliers", style='List Bullet')

    doc.add_heading("Risk Assessment", level=1)
    doc.add_paragraph("""
    Key risks identified for Q1 2025:
    1. Supply chain disruptions from geopolitical tensions (Medium risk)
    2. Labor market tightness in manufacturing regions (Low risk)
    3. Raw material price volatility (Medium risk)

    Mitigation strategies are in place for all identified risks.
    """)

    doc.save("operations_update.docx")
    print("Created operations_update.docx")


if __name__ == "__main__":
    print("Generating input files for P15 Multi-Format Integration...")
    print()

    create_sales_excel()
    create_marketing_pptx()
    create_finance_pdf()
    create_operations_docx()

    print()
    print("All files created successfully!")
